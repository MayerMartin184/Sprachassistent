"""Dokumente erstellen: Word (.docx), Excel (.xlsx), PowerPoint (.pptx) – als Dateien in den freigegebenen Ordnern.

Inhalte kommen als einfaches Markdown (Überschriften mit #, Aufzählungen mit -, Absätze) bzw. als Tabellenzeilen.
"""

from __future__ import annotations

import re
from typing import Any

from .base import Tool, schema
from .files import FileManager


def _blocks(markdown: str) -> list[tuple[str, str]]:
    """Zerlegt einfaches Markdown in (Typ, Text): h1/h2/h3, bullet, number, paragraph."""
    out: list[tuple[str, str]] = []
    for raw in markdown.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        if m := re.match(r"^(#{1,3})\s+(.*)", line):
            out.append((f"h{len(m.group(1))}", m.group(2).strip()))
        elif m := re.match(r"^\s*[-*•]\s+(.*)", line):
            out.append(("bullet", m.group(1).strip()))
        elif m := re.match(r"^\s*\d+[.)]\s+(.*)", line):
            out.append(("number", m.group(1).strip()))
        else:
            out.append(("paragraph", line.strip()))
    return out


def _clean(text: str) -> str:
    return re.sub(r"\*\*(.+?)\*\*|\*(.+?)\*|`(.+?)`", lambda m: m.group(1) or m.group(2) or m.group(3), text)


def create_docx(fm: FileManager, path: str, title: str, content: str, author: str | None = None) -> str:
    import docx
    from docx.shared import Pt

    target = fm.resolve(path if path.lower().endswith(".docx") else path + ".docx")
    if target.exists():
        raise FileExistsError(f"Datei existiert bereits: {fm.rel(target)}")
    target.parent.mkdir(parents=True, exist_ok=True)
    doc = docx.Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    if author:
        doc.core_properties.author = author
    doc.core_properties.title = title
    doc.add_heading(title, level=0)
    for kind, text in _blocks(content):
        text = _clean(text)
        if kind.startswith("h"):
            doc.add_heading(text, level=int(kind[1]))
        elif kind == "bullet":
            doc.add_paragraph(text, style="List Bullet")
        elif kind == "number":
            doc.add_paragraph(text, style="List Number")
        else:
            doc.add_paragraph(text)
    doc.save(str(target))
    return f"Word-Dokument erstellt: {fm.rel(target)}"


def create_xlsx(fm: FileManager, path: str, sheets: list[dict[str, Any]]) -> str:
    import openpyxl
    from openpyxl.styles import Font, PatternFill
    from openpyxl.utils import get_column_letter

    target = fm.resolve(path if path.lower().endswith(".xlsx") else path + ".xlsx")
    if target.exists():
        raise FileExistsError(f"Datei existiert bereits: {fm.rel(target)}")
    target.parent.mkdir(parents=True, exist_ok=True)
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for sheet in sheets:
        ws = wb.create_sheet(str(sheet.get("name") or "Tabelle")[:31])
        rows = sheet.get("rows") or []
        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row, start=1):
                ws.cell(row=r, column=c, value=_coerce(value))
        if rows and sheet.get("header", True):
            for cell in ws[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="1F4E5A")
            ws.freeze_panes = "A2"
        for c in range(1, (max(len(r) for r in rows) if rows else 0) + 1):
            width = max((len(str(row[c - 1])) for row in rows if len(row) >= c), default=8)
            ws.column_dimensions[get_column_letter(c)].width = min(max(10, width + 2), 60)
    if not wb.worksheets:
        wb.create_sheet("Tabelle")
    wb.save(str(target))
    return f"Excel-Datei erstellt: {fm.rel(target)} ({len(sheets)} Blatt/Blätter)"


def _coerce(value: Any) -> Any:
    if isinstance(value, str):
        v = value.strip()
        if re.fullmatch(r"-?\d+", v):
            return int(v)
        if re.fullmatch(r"-?\d+[.,]\d+", v):
            return float(v.replace(",", "."))
    return value


def create_pptx(fm: FileManager, path: str, title: str, subtitle: str | None, slides: list[dict[str, Any]]) -> str:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    target = fm.resolve(path if path.lower().endswith(".pptx") else path + ".pptx")
    if target.exists():
        raise FileExistsError(f"Datei existiert bereits: {fm.rel(target)}")
    target.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    dark, accent, light = RGBColor(0x0C, 0x17, 0x1B), RGBColor(0xA7, 0xE3, 0xEA), RGBColor(0xF2, 0xF7, 0xF8)

    def background(slide) -> None:  # noqa: ANN001
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = dark

    first = prs.slides.add_slide(prs.slide_layouts[6])
    background(first)
    box = first.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size, p.font.bold, p.font.color.rgb = Pt(44), True, light
    if subtitle:
        p2 = box.text_frame.add_paragraph()
        p2.text = subtitle
        p2.font.size, p2.font.color.rgb = Pt(20), accent

    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background(slide)
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tp = tb.text_frame.paragraphs[0]
        tp.text = str(s.get("title") or "")
        tp.font.size, tp.font.bold, tp.font.color.rgb = Pt(30), True, accent
        body = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.3))
        body.text_frame.word_wrap = True
        bullets = s.get("bullets") or []
        for i, bullet in enumerate(bullets):
            bp = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
            bp.text = "•  " + _clean(str(bullet))
            bp.font.size, bp.font.color.rgb = Pt(20), light
            bp.space_after = Pt(10)
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = str(s["notes"])
    prs.save(str(target))
    return f"Präsentation erstellt: {fm.rel(target)} ({len(slides) + 1} Folien)"


def build_tools(fm: FileManager, author: str | None = None) -> list[Tool]:
    hint = "Pfad wie bei den Dateiwerkzeugen, z. B. 'Dokumente/Angebote/Angebot_Popescu'."
    return [
        Tool(
            "create_docx",
            f"Erstellt ein Word-Dokument aus einfachem Markdown (#, ##, ### Überschriften; - Aufzählung; 1. Nummerierung; Absätze). {hint}",
            schema({"path": {"type": "string"}, "title": {"type": "string"}, "content": {"type": "string"}}, ["path", "title", "content"]),
            lambda path, title, content: create_docx(fm, path, title, content, author),
        ),
        Tool(
            "create_xlsx",
            f"Erstellt eine Excel-Datei. sheets: Liste von {{name, rows}}; rows ist eine Liste von Zeilen (Listen), erste Zeile = Überschriften. Zahlen als Zahlen. {hint}",
            schema(
                {
                    "path": {"type": "string"},
                    "sheets": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "name": {"type": "string"},
                                "rows": {"type": "array", "items": {"type": "array", "items": {}}},
                                "header": {"type": "boolean"},
                            },
                            "required": ["name", "rows"],
                            "additionalProperties": False,
                        },
                    },
                },
                ["path", "sheets"],
            ),
            lambda path, sheets: create_xlsx(fm, path, sheets),
        ),
        Tool(
            "create_pptx",
            f"Erstellt eine PowerPoint-Präsentation im dunklen Firmenstil: Titelfolie plus Folien mit Titel und Stichpunkten (optional Notizen). {hint}",
            schema(
                {
                    "path": {"type": "string"},
                    "title": {"type": "string"},
                    "subtitle": {"type": "string"},
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "bullets": {"type": "array", "items": {"type": "string"}},
                                "notes": {"type": "string"},
                            },
                            "required": ["title", "bullets"],
                            "additionalProperties": False,
                        },
                    },
                },
                ["path", "title", "slides"],
            ),
            lambda path, title, slides, subtitle=None: create_pptx(fm, path, title, subtitle, slides),
        ),
    ]
